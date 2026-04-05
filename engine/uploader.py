"""
uploader.py
-----------
File Upload & Text Extraction Module
Supports: PDF (pdfplumber), PPTX (python-pptx)
Extracts text → splits into paragraphs/slides → preprocesses → adds to live index
"""

import os
import re
import json
from .preprocessor import Preprocessor

try:
    import pdfplumber
    PDF_OK = True
except ImportError:
    PDF_OK = False

try:
    from pptx import Presentation
    PPTX_OK = True
except ImportError:
    PPTX_OK = False


class Uploader:
    """
    Handles PDF and PPTX uploads.
    Extracts text, detects topic/section boundaries,
    creates document objects, and injects them into the live index.
    """

    def __init__(self, indexer):
        self.indexer = indexer
        self.preprocessor = Preprocessor()
        self.upload_folder = 'uploads'
        os.makedirs(self.upload_folder, exist_ok=True)

        # Track uploaded files
        self.uploaded_files = []

    # ------------------------------------------------------------------ #
    #  PDF TEXT EXTRACTION                                                 #
    # ------------------------------------------------------------------ #

    def extract_pdf(self, filepath):
        """
        Extracts text from every page of a PDF using pdfplumber.
        Returns list of { page_num, text } dicts.
        """
        if not PDF_OK:
            raise ImportError("pdfplumber not installed. Run: pip install pdfplumber")

        pages = []
        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages):
                raw = page.extract_text()
                if raw and raw.strip():
                    pages.append({
                        'page_num': i + 1,
                        'text': raw.strip()
                    })
        return pages

    # ------------------------------------------------------------------ #
    #  PPTX TEXT EXTRACTION                                                #
    # ------------------------------------------------------------------ #

    def extract_pptx(self, filepath):
        """
        Extracts text from every slide of a PPTX file.
        Returns list of { slide_num, title, text } dicts.
        """
        if not PPTX_OK:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")

        prs = Presentation(filepath)
        slides = []
        for i, slide in enumerate(prs.slides):
            title_text = ''
            body_texts = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                # Detect title placeholder
                if shape.shape_type == 13 or (hasattr(shape, 'placeholder_format') and
                        shape.placeholder_format and shape.placeholder_format.idx == 0):
                    title_text = text
                else:
                    body_texts.append(text)

            combined = ' '.join(body_texts)
            if combined.strip() or title_text.strip():
                slides.append({
                    'slide_num': i + 1,
                    'title': title_text or f'Slide {i + 1}',
                    'text': combined.strip()
                })
        return slides

    # ------------------------------------------------------------------ #
    #  SMART PARAGRAPH SPLITTING                                           #
    # ------------------------------------------------------------------ #

    def split_into_paragraphs(self, text, min_words=20):
        """
        Splits a long text block into meaningful paragraphs.
        Splits on double newlines or sentence boundaries.
        Filters out very short fragments.
        """
        # Split on double newlines first
        raw_chunks = re.split(r'\n{2,}', text)
        paragraphs = []
        for chunk in raw_chunks:
            chunk = chunk.strip().replace('\n', ' ')
            chunk = re.sub(r'\s+', ' ', chunk)
            if len(chunk.split()) >= min_words:
                paragraphs.append(chunk)

        # If no good splits found, return as single block
        if not paragraphs and len(text.split()) >= min_words:
            return [text.strip()]
        return paragraphs

    def detect_topic(self, text, fallback='Uploaded Content'):
        """
        Tries to detect a topic/heading from the first line of text.
        Falls back to the fallback label if nothing useful found.
        """
        lines = text.strip().split('\n')
        for line in lines[:3]:
            line = line.strip()
            # A good topic line: not too long, not too short, mostly letters
            if 5 < len(line) < 80 and sum(c.isalpha() for c in line) / max(len(line), 1) > 0.5:
                return line
        # Extract first meaningful sentence
        sentences = re.split(r'[.!?]', text)
        if sentences and len(sentences[0].split()) <= 12:
            return sentences[0].strip()
        return fallback

    def detect_unit(self, text):
        """
        Attempts to detect IRT unit from content keywords.
        Falls back to 'Uploaded' if no unit detected.
        """
        text_lower = text.lower()
        unit_keywords = {
            'Unit 1': ['information retrieval', 'ir system', 'data retrieval', 'search architecture',
                       'objectives of ir', 'definition of ir'],
            'Unit 2': ['boolean model', 'tf-idf', 'vector space', 'probabilistic', 'language model',
                       'precision', 'recall', 'f-measure', 'bm25', 'cosine similarity'],
            'Unit 3': ['inverted index', 'postings', 'classification', 'clustering', 'k-means',
                       'naive bayes', 'svm', 'knn', 'k-nearest', 'hierarchical'],
            'Unit 4': ['pagerank', 'web crawler', 'search engine', 'hyperlink', 'link analysis', 'seo'],
            'Unit 5': ['stemming', 'lemmatization', 'tokenization', 'stopword', 'n-gram',
                       'sentiment', 'nlp', 'recommendation', 'document scoring']
        }
        scores = {unit: 0 for unit in unit_keywords}
        for unit, keywords in unit_keywords.items():
            for kw in keywords:
                if kw in text_lower:
                    scores[unit] += 1
        best_unit = max(scores, key=scores.get)
        return best_unit if scores[best_unit] > 0 else 'Uploaded'

    # ------------------------------------------------------------------ #
    #  MAIN PROCESSING PIPELINE                                            #
    # ------------------------------------------------------------------ #

    def process_file(self, filepath, filename, subject='Uploaded Notes'):
        """
        Full pipeline:
        1. Extract text (PDF or PPTX)
        2. Split into sections/paragraphs
        3. Detect topic and unit for each section
        4. Preprocess text (tokenize, stem, lemmatize, remove stopwords)
        5. Add to live index dynamically
        Returns: list of created document objects
        """
        ext = os.path.splitext(filename)[1].lower()
        created_docs = []
        base_id = max((d['id'] for d in self.indexer.documents), default=100) + 1

        if ext == '.pdf':
            pages = self.extract_pdf(filepath)
            for page in pages:
                paragraphs = self.split_into_paragraphs(page['text'])
                if not paragraphs:
                    paragraphs = [page['text']]
                for j, para in enumerate(paragraphs):
                    topic = self.detect_topic(para, fallback=f'{filename} — Page {page["page_num"]}')
                    unit = self.detect_unit(para)
                    doc = {
                        'id': base_id,
                        'subject': subject,
                        'unit': unit,
                        'topic': topic,
                        'content': para,
                        'source': filename,
                        'source_type': 'pdf',
                        'page': page['page_num']
                    }
                    self.indexer.add_document(doc)
                    created_docs.append(doc)
                    base_id += 1

        elif ext in ('.pptx', '.ppt'):
            slides = self.extract_pptx(filepath)
            for slide in slides:
                content = slide['text'] if slide['text'] else slide['title']
                if not content.strip():
                    continue
                unit = self.detect_unit(content + ' ' + slide['title'])
                doc = {
                    'id': base_id,
                    'subject': subject,
                    'unit': unit,
                    'topic': slide['title'] or f'{filename} — Slide {slide["slide_num"]}',
                    'content': content,
                    'source': filename,
                    'source_type': 'pptx',
                    'slide': slide['slide_num']
                }
                self.indexer.add_document(doc)
                created_docs.append(doc)
                base_id += 1

        else:
            raise ValueError(f"Unsupported file type: {ext}. Only PDF and PPTX are supported.")

        # Record upload
        self.uploaded_files.append({
            'filename': filename,
            'type': ext,
            'sections_extracted': len(created_docs)
        })

        print(f"[Uploader] Processed '{filename}': {len(created_docs)} sections added to index.")
        return created_docs

    def get_upload_history(self):
        return self.uploaded_files

    def get_uploaded_docs(self):
        """Returns only the documents that came from uploaded files."""
        return [d for d in self.indexer.documents if d.get('source_type') in ('pdf', 'pptx')]